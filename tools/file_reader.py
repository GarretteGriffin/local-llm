"""
File Reader Tool - Handles all document types.
Supports: Office (docx, xlsx, pptx), PDF, text files, CSV, JSON, XML
"""
import os
from pathlib import Path
from typing import Optional, Dict, Any, Union
from dataclasses import dataclass
import tempfile
import json

from config import settings


@dataclass
class FileContent:
    """Extracted file content"""
    filename: str
    file_type: str
    content: str
    metadata: Dict[str, Any]
    page_count: Optional[int] = None
    error: Optional[str] = None


class FileReaderTool:
    """
    Reads and extracts content from various file types.
    All processing is automatic - user just uploads files.
    """
    
    SUPPORTED_TYPES = {
        # Office
        '.docx': 'word',
        '.doc': 'word_legacy',
        '.xlsx': 'excel',
        '.xls': 'excel_legacy',
        '.pptx': 'powerpoint',
        '.ppt': 'powerpoint_legacy',
        # PDF
        '.pdf': 'pdf',
        # Data
        '.csv': 'csv',
        '.json': 'json',
        '.xml': 'xml',
        # Text
        '.txt': 'text',
        '.md': 'markdown',
        '.rtf': 'rtf',
    }
    
    def __init__(self):
        self.settings = settings
        self.max_size = settings.max_file_size_mb * 1024 * 1024
    
    def read(self, file_path: Union[str, Path], file_bytes: Optional[bytes] = None) -> FileContent:
        """
        Read content from a file.
        
        Args:
            file_path: Path to the file or filename
            file_bytes: Optional file bytes (for uploaded files)
            
        Returns:
            FileContent with extracted text and metadata
        """
        path = Path(file_path)
        extension = path.suffix.lower()
        filename = path.name
        
        if extension not in self.SUPPORTED_TYPES:
            return FileContent(
                filename=filename,
                file_type=extension,
                content="",
                metadata={},
                error=f"Unsupported file type: {extension}"
            )
        
        file_type = self.SUPPORTED_TYPES[extension]
        
        # Handle file bytes (uploaded files)
        if file_bytes:
            if len(file_bytes) > self.max_size:
                return FileContent(
                    filename=filename,
                    file_type=file_type,
                    content="",
                    metadata={},
                    error=f"File too large (max {self.settings.max_file_size_mb}MB)"
                )
            # Write to temp file
            with tempfile.NamedTemporaryFile(delete=False, suffix=extension) as tmp:
                tmp.write(file_bytes)
                temp_path = tmp.name
            try:
                result = self._read_by_type(temp_path, file_type, filename)
            finally:
                os.unlink(temp_path)
            return result
        
        # Check file exists
        if not path.exists():
            return FileContent(
                filename=filename,
                file_type=file_type,
                content="",
                metadata={},
                error=f"File not found: {file_path}"
            )
        
        # Check file size
        if path.stat().st_size > self.max_size:
            return FileContent(
                filename=filename,
                file_type=file_type,
                content="",
                metadata={},
                error=f"File too large (max {self.settings.max_file_size_mb}MB)"
            )
        
        return self._read_by_type(str(path), file_type, filename)
    
    def _read_by_type(self, file_path: str, file_type: str, filename: str) -> FileContent:
        """Read file based on type"""
        try:
            if file_type == 'word':
                return self._read_docx(file_path, filename)
            elif file_type == 'excel':
                return self._read_xlsx(file_path, filename)
            elif file_type == 'powerpoint':
                return self._read_pptx(file_path, filename)
            elif file_type == 'pdf':
                return self._read_pdf(file_path, filename)
            elif file_type == 'csv':
                return self._read_csv(file_path, filename)
            elif file_type == 'json':
                return self._read_json(file_path, filename)
            elif file_type == 'xml':
                return self._read_xml(file_path, filename)
            elif file_type in ('text', 'markdown'):
                return self._read_text(file_path, filename, file_type)
            else:
                return FileContent(
                    filename=filename,
                    file_type=file_type,
                    content="",
                    metadata={},
                    error=f"Reader not implemented for: {file_type}"
                )
        except Exception as e:
            return FileContent(
                filename=filename,
                file_type=file_type,
                content="",
                metadata={},
                error=f"Error reading file: {str(e)}"
            )
    
    def _read_docx(self, file_path: str, filename: str) -> FileContent:
        """Read Word documents"""
        try:
            from docx import Document
        except ImportError:
            return FileContent(
                filename=filename, file_type='word', content="",
                metadata={}, error="python-docx not installed"
            )
        
        doc = Document(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        
        # Also extract tables
        tables_text = []
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                tables_text.append(" | ".join(cells))
        
        content = "\n\n".join(paragraphs)
        if tables_text:
            content += "\n\n[Tables]\n" + "\n".join(tables_text)
        
        return FileContent(
            filename=filename,
            file_type='word',
            content=content,
            metadata={
                'paragraphs': len(paragraphs),
                'tables': len(doc.tables)
            }
        )
    
    def _read_xlsx(self, file_path: str, filename: str) -> FileContent:
        """Read Excel spreadsheets - optimized for 32k context window"""
        try:
            from openpyxl import load_workbook
        except ImportError:
            return FileContent(
                filename=filename, file_type='excel', content="",
                metadata={}, error="openpyxl not installed"
            )
        
        # 32k context = ~120k chars max, but leave room for prompt/response
        MAX_ROWS_PER_SHEET = 1000  # Reasonable sample
        MAX_TOTAL_CHARS = 80000   # ~20k tokens - fits well in 32k context
        
        wb = load_workbook(file_path, data_only=True)
        sheets_content = []
        total_rows_all_sheets = 0
        total_chars = 0
        
        for sheet_name in wb.sheetnames:
            if total_chars >= MAX_TOTAL_CHARS:
                sheets_content.append(f"\n[Additional sheets truncated - context limit reached]")
                break
                
            sheet = wb[sheet_name]
            all_rows = []
            
            for row in sheet.iter_rows(values_only=True):
                if any(cell is not None for cell in row):
                    row_str = " | ".join(str(cell) if cell is not None else "" for cell in row)
                    all_rows.append(row_str)
            
            total_rows_all_sheets += len(all_rows)
            
            if not all_rows:
                continue
            
            # Include as many rows as possible
            if len(all_rows) > MAX_ROWS_PER_SHEET:
                # Keep header + evenly sampled data
                header = all_rows[0]
                data_rows = all_rows[1:]
                step = max(1, len(data_rows) // (MAX_ROWS_PER_SHEET - 1))
                sampled = [data_rows[i] for i in range(0, len(data_rows), step)][:MAX_ROWS_PER_SHEET-1]
                
                sheet_text = f"[Sheet: {sheet_name}] ({len(all_rows):,} total rows, {len(sampled)+1:,} included)\n"
                sheet_text += header + "\n" + "\n".join(sampled)
            else:
                sheet_text = f"[Sheet: {sheet_name}] ({len(all_rows):,} rows - complete)\n" + "\n".join(all_rows)
            
            total_chars += len(sheet_text)
            sheets_content.append(sheet_text)
        
        full_content = "\n\n".join(sheets_content)
        
        # Final check
        if len(full_content) > MAX_TOTAL_CHARS:
            full_content = full_content[:MAX_TOTAL_CHARS] + f"\n\n[Content truncated at {MAX_TOTAL_CHARS:,} chars. Total: {total_rows_all_sheets:,} rows across {len(wb.sheetnames)} sheets]"
        
        return FileContent(
            filename=filename,
            file_type='excel',
            content=full_content,
            metadata={
                'sheets': wb.sheetnames,
                'sheet_count': len(wb.sheetnames),
                'total_rows': total_rows_all_sheets,
                'content_chars': len(full_content)
            }
        )
    
    def _read_pptx(self, file_path: str, filename: str) -> FileContent:
        """Read PowerPoint presentations"""
        try:
            from pptx import Presentation
        except ImportError:
            return FileContent(
                filename=filename, file_type='powerpoint', content="",
                metadata={}, error="python-pptx not installed"
            )
        
        prs = Presentation(file_path)
        slides_content = []
        
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
            
            if texts:
                slides_content.append(f"[Slide {i}]\n" + "\n".join(texts))
        
        return FileContent(
            filename=filename,
            file_type='powerpoint',
            content="\n\n".join(slides_content),
            metadata={'slide_count': len(prs.slides)},
            page_count=len(prs.slides)
        )
    
    def _read_pdf(self, file_path: str, filename: str) -> FileContent:
        """Read PDF documents"""
        # Try PyMuPDF first (faster, better quality)
        try:
            import fitz  # PyMuPDF
            
            doc = fitz.open(file_path)
            pages_content = []
            
            for i, page in enumerate(doc, 1):
                text = page.get_text()
                if text.strip():
                    pages_content.append(f"[Page {i}]\n{text}")
            
            return FileContent(
                filename=filename,
                file_type='pdf',
                content="\n\n".join(pages_content),
                metadata={'pages': len(doc)},
                page_count=len(doc)
            )
            
        except ImportError:
            pass
        
        # Fallback to pdfplumber
        try:
            import pdfplumber
            
            with pdfplumber.open(file_path) as pdf:
                pages_content = []
                for i, page in enumerate(pdf.pages, 1):
                    text = page.extract_text()
                    if text and text.strip():
                        pages_content.append(f"[Page {i}]\n{text}")
                
                return FileContent(
                    filename=filename,
                    file_type='pdf',
                    content="\n\n".join(pages_content),
                    metadata={'pages': len(pdf.pages)},
                    page_count=len(pdf.pages)
                )
                
        except ImportError:
            return FileContent(
                filename=filename, file_type='pdf', content="",
                metadata={}, error="No PDF library installed (install PyMuPDF or pdfplumber)"
            )
    
    def _read_csv(self, file_path: str, filename: str) -> FileContent:
        """Read CSV files"""
        import csv
        
        rows = []
        with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
            reader = csv.reader(f)
            for i, row in enumerate(reader):
                if i < 1000:  # Limit rows
                    rows.append(" | ".join(row))
                else:
                    rows.append(f"... ({i} rows total, showing first 1000)")
                    break
        
        return FileContent(
            filename=filename,
            file_type='csv',
            content="\n".join(rows),
            metadata={'rows': len(rows)}
        )
    
    def _read_json(self, file_path: str, filename: str) -> FileContent:
        """Read JSON files"""
        with open(file_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        
        # Pretty print with limit
        content = json.dumps(data, indent=2)
        if len(content) > 50000:
            content = content[:50000] + "\n... (truncated)"
        
        return FileContent(
            filename=filename,
            file_type='json',
            content=content,
            metadata={'type': type(data).__name__}
        )
    
    def _read_xml(self, file_path: str, filename: str) -> FileContent:
        """Read XML files"""
        with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
            content = f.read()
        
        if len(content) > 50000:
            content = content[:50000] + "\n... (truncated)"
        
        return FileContent(
            filename=filename,
            file_type='xml',
            content=content,
            metadata={}
        )
    
    def _read_text(self, file_path: str, filename: str, file_type: str) -> FileContent:
        """Read text files"""
        with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
            content = f.read()
        
        if len(content) > 100000:
            content = content[:100000] + "\n... (truncated)"
        
        return FileContent(
            filename=filename,
            file_type=file_type,
            content=content,
            metadata={'chars': len(content)}
        )
